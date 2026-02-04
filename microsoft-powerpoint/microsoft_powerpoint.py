from autohive_integrations_sdk import (
    Integration, ExecutionContext, ActionHandler, ActionResult
)
from typing import Dict, Any, Optional, Tuple
import io
import re
import asyncio

microsoft_powerpoint = Integration.load()

GRAPH_API_BASE_URL = "https://graph.microsoft.com/v1.0"

THUMBNAIL_SIZES = {
    "small": {"width": 176, "height": 144},
    "medium": {"width": 800, "height": 600},
    "large": {"width": 1600, "height": 1200}
}

MAX_SIMPLE_UPLOAD_SIZE = 4 * 1024 * 1024
MAX_DOWNLOAD_SIZE = 50 * 1024 * 1024  # 50MB limit for downloads
COPY_POLL_INTERVAL = 1.0  # seconds
COPY_POLL_MAX_ATTEMPTS = 30  # 30 seconds max wait


def odata_escape(value: str) -> str:
    """Escape single quotes for OData filter strings."""
    return value.replace("'", "''")


def validate_path(path: str) -> bool:
    """Validate that a path doesn't contain path traversal sequences."""
    if not path:
        return True
    # Check for path traversal attempts
    if '..' in path or path.startswith('/..') or '/../' in path:
        return False
    return True


async def get_file_metadata(context: ExecutionContext, item_id: str) -> Dict[str, Any]:
    """Get file metadata including size and eTag."""
    response = await context.fetch(
        f"{GRAPH_API_BASE_URL}/me/drive/items/{item_id}",
        method="GET",
        params={"$select": "id,name,size,eTag,cTag"}
    )
    return response


async def download_file_content(context: ExecutionContext, item_id: str, check_size: bool = True) -> bytes:
    """Download file content with optional size check to prevent memory issues."""
    if check_size:
        metadata = await get_file_metadata(context, item_id)
        file_size = metadata.get('size', 0)
        if file_size > MAX_DOWNLOAD_SIZE:
            raise ValueError(f"File size ({file_size} bytes) exceeds download limit ({MAX_DOWNLOAD_SIZE} bytes). File is too large to process.")

    response = await context.fetch(
        f"{GRAPH_API_BASE_URL}/me/drive/items/{item_id}/content",
        method="GET",
        raw_response=True
    )
    return response


async def get_slide_count_from_metadata(context: ExecutionContext, item_id: str) -> Optional[int]:
    """Try to get slide count without downloading the file. Returns None if not available."""
    # Microsoft Graph doesn't expose slide count directly, so we return None
    # This is a placeholder for future optimization if MS adds this capability
    return None


async def overwrite_file_content(context: ExecutionContext, item_id: str, content: bytes) -> Dict[str, Any]:
    """Overwrite file content by item ID (simpler and more reliable than path-based upload)."""
    if len(content) > MAX_SIMPLE_UPLOAD_SIZE:
        raise ValueError(f"File size ({len(content)} bytes) exceeds simple upload limit ({MAX_SIMPLE_UPLOAD_SIZE} bytes). Large file upload not yet supported.")
    
    response = await context.fetch(
        f"{GRAPH_API_BASE_URL}/me/drive/items/{item_id}/content",
        method="PUT",
        headers={
            "Content-Type": "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        },
        data=content
    )
    return response


async def upload_file_content(context: ExecutionContext, folder_path: str, file_name: str, content: bytes) -> Dict[str, Any]:
    """Upload new file content to a folder path."""
    if len(content) > MAX_SIMPLE_UPLOAD_SIZE:
        raise ValueError(f"File size ({len(content)} bytes) exceeds simple upload limit ({MAX_SIMPLE_UPLOAD_SIZE} bytes). Large file upload not yet supported.")
    
    if folder_path and folder_path != "/":
        upload_url = f"{GRAPH_API_BASE_URL}/me/drive/root:/{folder_path.strip('/')}/{file_name}:/content"
    else:
        upload_url = f"{GRAPH_API_BASE_URL}/me/drive/root:/{file_name}:/content"
    
    response = await context.fetch(
        upload_url,
        method="PUT",
        headers={
            "Content-Type": "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        },
        data=content
    )
    return response


def create_blank_pptx() -> bytes:
    try:
        from pptx import Presentation
        prs = Presentation()
        buffer = io.BytesIO()
        try:
            prs.save(buffer)
            buffer.seek(0)
            return buffer.read()
        finally:
            buffer.close()
    except ImportError:
        raise ImportError("python-pptx library is required for creating presentations. Install with: pip install python-pptx")


@microsoft_powerpoint.action("powerpoint_list_presentations")
class ListPresentationsAction(ActionHandler):
    """Find accessible PowerPoint presentations (.pptx) in OneDrive/SharePoint."""

    async def execute(self, inputs: Dict[str, Any], context: ExecutionContext) -> ActionResult:
        try:
            folder_path = inputs.get('folder_path', '')
            name_contains = inputs.get('name_contains', '')
            page_size = min(inputs.get('page_size', 25), 100)
            page_token = inputs.get('page_token')

            if folder_path and folder_path != "/":
                base_url = f"{GRAPH_API_BASE_URL}/me/drive/root:/{folder_path.strip('/')}:/children"
            else:
                base_url = f"{GRAPH_API_BASE_URL}/me/drive/root/children"

            filter_parts = ["endswith(name, '.pptx')"]
            if name_contains:
                filter_parts.append(f"contains(name, '{odata_escape(name_contains)}')")

            params = {
                "$filter": " and ".join(filter_parts),
                "$top": page_size,
                "$select": "id,name,webUrl,lastModifiedDateTime,size,createdDateTime"
            }

            if page_token:
                response = await context.fetch(page_token, method="GET")
            else:
                response = await context.fetch(base_url, method="GET", params=params)

            items = response.get('value', [])
            presentations = [
                {
                    "id": item.get('id'),
                    "name": item.get('name'),
                    "webUrl": item.get('webUrl'),
                    "lastModifiedDateTime": item.get('lastModifiedDateTime'),
                    "size": item.get('size')
                }
                for item in items
            ]

            next_page_token = response.get('@odata.nextLink')

            return ActionResult(
                data={
                    "presentations": presentations,
                    "next_page_token": next_page_token,
                    "result": True
                },
                cost_usd=0.0
            )

        except Exception as e:
            return ActionResult(
                data={"presentations": [], "result": False, "error": str(e)},
                cost_usd=0.0
            )


@microsoft_powerpoint.action("powerpoint_get_presentation")
class GetPresentationAction(ActionHandler):
    """Retrieve presentation properties including file info, author, and last modified details."""

    async def execute(self, inputs: Dict[str, Any], context: ExecutionContext) -> ActionResult:
        try:
            presentation_id = inputs['presentation_id']

            response = await context.fetch(
                f"{GRAPH_API_BASE_URL}/me/drive/items/{presentation_id}",
                method="GET",
                params={
                    "$select": "id,name,size,webUrl,createdDateTime,lastModifiedDateTime,createdBy,lastModifiedBy,parentReference"
                }
            )

            return ActionResult(
                data={
                    "id": response.get('id'),
                    "name": response.get('name'),
                    "size": response.get('size'),
                    "webUrl": response.get('webUrl'),
                    "createdDateTime": response.get('createdDateTime'),
                    "lastModifiedDateTime": response.get('lastModifiedDateTime'),
                    "createdBy": response.get('createdBy'),
                    "lastModifiedBy": response.get('lastModifiedBy'),
                    "parentReference": response.get('parentReference'),
                    "result": True
                },
                cost_usd=0.0
            )

        except Exception as e:
            return ActionResult(
                data={"result": False, "error": str(e)},
                cost_usd=0.0
            )


@microsoft_powerpoint.action("powerpoint_get_slides")
class GetSlidesAction(ActionHandler):
    """List all slides in a presentation with their basic metadata."""

    async def execute(self, inputs: Dict[str, Any], context: ExecutionContext) -> ActionResult:
        try:
            from pptx import Presentation

            presentation_id = inputs['presentation_id']
            include_thumbnails = inputs.get('include_thumbnails', True)
            thumbnail_size = inputs.get('thumbnail_size', 'medium')

            content_bytes = await download_file_content(context, presentation_id)
            prs = Presentation(io.BytesIO(content_bytes))

            slides = []
            thumbnail_data = {}

            if include_thumbnails:
                try:
                    thumbnails_response = await context.fetch(
                        f"{GRAPH_API_BASE_URL}/me/drive/items/{presentation_id}/thumbnails",
                        method="GET"
                    )
                    thumbnail_sets = thumbnails_response.get('value', [])
                    if thumbnail_sets:
                        size_data = thumbnail_sets[0].get(thumbnail_size, thumbnail_sets[0].get('medium', {}))
                        thumbnail_data = {
                            "url": size_data.get('url'),
                            "width": size_data.get('width'),
                            "height": size_data.get('height')
                        }
                except Exception:
                    pass  # Thumbnail fetch is non-critical - continue without thumbnails

            for i, slide in enumerate(prs.slides):
                slide_info = {
                    "index": i + 1,
                    "id": str(slide.slide_id)
                }

                if include_thumbnails and thumbnail_data.get('url'):
                    slide_info["thumbnailUrl"] = thumbnail_data.get('url')
                    slide_info["thumbnailWidth"] = thumbnail_data.get('width')
                    slide_info["thumbnailHeight"] = thumbnail_data.get('height')
                    slide_info["thumbnailNote"] = "Thumbnail is for entire presentation, not individual slide"

                slides.append(slide_info)

            return ActionResult(
                data={
                    "slides": slides,
                    "slide_count": len(slides),
                    "result": True
                },
                cost_usd=0.0
            )

        except ImportError:
            return ActionResult(
                data={"slides": [], "slide_count": 0, "result": False, "error": "python-pptx library is required. Install with: pip install python-pptx"},
                cost_usd=0.0
            )
        except Exception as e:
            return ActionResult(
                data={"slides": [], "slide_count": 0, "result": False, "error": str(e)},
                cost_usd=0.0
            )


@microsoft_powerpoint.action("powerpoint_get_slide")
class GetSlideAction(ActionHandler):
    """Get details for a specific slide by its index (1-based)."""

    async def execute(self, inputs: Dict[str, Any], context: ExecutionContext) -> ActionResult:
        try:
            from pptx import Presentation

            presentation_id = inputs['presentation_id']
            slide_index = inputs['slide_index']
            include_thumbnail = inputs.get('include_thumbnail', True)
            thumbnail_size = inputs.get('thumbnail_size', 'large')

            if slide_index < 1:
                return ActionResult(
                    data={"result": False, "error": "Slide index must be 1 or greater"},
                    cost_usd=0.0
                )

            content_bytes = await download_file_content(context, presentation_id)
            prs = Presentation(io.BytesIO(content_bytes))
            
            if slide_index > len(prs.slides):
                return ActionResult(
                    data={"result": False, "error": f"Slide index {slide_index} is out of range. Presentation has {len(prs.slides)} slides."},
                    cost_usd=0.0
                )

            slide = prs.slides[slide_index - 1]
            
            result_data = {
                "index": slide_index,
                "id": str(slide.slide_id),
                "result": True
            }

            if include_thumbnail:
                try:
                    thumbnails_response = await context.fetch(
                        f"{GRAPH_API_BASE_URL}/me/drive/items/{presentation_id}/thumbnails",
                        method="GET"
                    )
                    thumbnail_sets = thumbnails_response.get('value', [])
                    if thumbnail_sets:
                        size_data = thumbnail_sets[0].get(thumbnail_size, thumbnail_sets[0].get('large', {}))
                        result_data["thumbnailUrl"] = size_data.get('url')
                        result_data["thumbnailWidth"] = size_data.get('width')
                        result_data["thumbnailHeight"] = size_data.get('height')
                        result_data["thumbnailNote"] = "Thumbnail is for entire presentation, not individual slide"
                except Exception:
                    pass  # Thumbnail fetch is non-critical - continue without thumbnail

            return ActionResult(
                data=result_data,
                cost_usd=0.0
            )

        except ImportError:
            return ActionResult(
                data={"result": False, "error": "python-pptx library is required. Install with: pip install python-pptx"},
                cost_usd=0.0
            )
        except Exception as e:
            return ActionResult(
                data={"result": False, "error": str(e)},
                cost_usd=0.0
            )


@microsoft_powerpoint.action("powerpoint_create_presentation")
class CreatePresentationAction(ActionHandler):
    """Create a new PowerPoint presentation in the specified folder."""

    async def execute(self, inputs: Dict[str, Any], context: ExecutionContext) -> ActionResult:
        try:
            name = inputs['name']
            folder_path = inputs.get('folder_path', '')
            template_id = inputs.get('template_id')

            # Validate paths for security
            if not validate_path(name):
                return ActionResult(
                    data={"result": False, "error": "Invalid presentation name: path traversal not allowed"},
                    cost_usd=0.0
                )
            if folder_path and not validate_path(folder_path):
                return ActionResult(
                    data={"result": False, "error": "Invalid folder path: path traversal not allowed"},
                    cost_usd=0.0
                )

            if not name.endswith('.pptx'):
                file_name = f"{name}.pptx"
            else:
                file_name = name

            if template_id:
                if folder_path and folder_path != "/":
                    dest_path = f"/{folder_path.strip('/')}"
                else:
                    dest_path = "/"

                copy_body = {
                    "name": file_name,
                    "parentReference": {
                        "path": f"/drive/root:{dest_path}" if dest_path != "/" else "/drive/root:"
                    }
                }

                # The copy endpoint returns 202 Accepted with a Location header
                # containing a monitor URL, not the copied item directly
                response = await context.fetch(
                    f"{GRAPH_API_BASE_URL}/me/drive/items/{template_id}/copy",
                    method="POST",
                    json=copy_body,
                    include_headers=True  # Request headers in response
                )

                # Check if we got a direct response (some cases return the item)
                if isinstance(response, dict) and response.get('id'):
                    return ActionResult(
                        data={
                            "id": response.get('id'),
                            "name": response.get('name'),
                            "webUrl": response.get('webUrl'),
                            "createdDateTime": response.get('createdDateTime'),
                            "result": True
                        },
                        cost_usd=0.0
                    )

                # For async operations, the copy was initiated
                # The file will be available shortly
                return ActionResult(
                    data={
                        "result": True,
                        "message": "Copy operation initiated. The presentation will be available shortly. Use list_presentations to find it.",
                        "name": file_name,
                        "async_operation": True
                    },
                    cost_usd=0.0
                )
            else:
                pptx_content = create_blank_pptx()
                response = await upload_file_content(context, folder_path, file_name, pptx_content)

                return ActionResult(
                    data={
                        "id": response.get('id'),
                        "name": response.get('name'),
                        "webUrl": response.get('webUrl'),
                        "createdDateTime": response.get('createdDateTime'),
                        "result": True
                    },
                    cost_usd=0.0
                )

        except Exception as e:
            return ActionResult(
                data={"result": False, "error": str(e)},
                cost_usd=0.0
            )


@microsoft_powerpoint.action("powerpoint_add_slide")
class AddSlideAction(ActionHandler):
    """Add a new slide to an existing presentation."""

    async def execute(self, inputs: Dict[str, Any], context: ExecutionContext) -> ActionResult:
        try:
            from pptx import Presentation

            presentation_id = inputs['presentation_id']
            position = inputs.get('position')
            layout = inputs.get('layout', 'blank')
            title = inputs.get('title')
            content = inputs.get('content')

            # Validate position if provided
            if position is not None and position < 1:
                return ActionResult(
                    data={"result": False, "error": "Position must be 1 or greater"},
                    cost_usd=0.0
                )

            content_bytes = await download_file_content(context, presentation_id)

            prs = Presentation(io.BytesIO(content_bytes))

            # Validate position against current slide count
            if position is not None and position > len(prs.slides) + 1:
                return ActionResult(
                    data={"result": False, "error": f"Position {position} is invalid. Presentation has {len(prs.slides)} slides, so position must be between 1 and {len(prs.slides) + 1}"},
                    cost_usd=0.0
                )

            layout_mapping = {
                'blank': 6,
                'title': 0,
                'titleContent': 1,
                'twoContent': 3,
                'comparison': 4,
                'titleOnly': 5,
                'contentCaption': 7,
                'pictureCaption': 8
            }

            layout_index = layout_mapping.get(layout, 6)
            if layout_index >= len(prs.slide_layouts):
                layout_index = 6 if 6 < len(prs.slide_layouts) else 0

            slide_layout = prs.slide_layouts[layout_index]

            if position is not None and position <= len(prs.slides):
                new_slide = prs.slides.add_slide(slide_layout)

                xml_slides = prs.slides._sldIdLst
                slides_list = list(xml_slides)
                last_slide = slides_list[-1]
                xml_slides.remove(last_slide)
                xml_slides.insert(position - 1, last_slide)
                slide_index = position
            else:
                new_slide = prs.slides.add_slide(slide_layout)
                slide_index = len(prs.slides)

            if title:
                from pptx.enum.shapes import PP_PLACEHOLDER
                title_types = {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE}
                for shape in new_slide.shapes:
                    if shape.has_text_frame and shape.is_placeholder:
                        if shape.placeholder_format.type in title_types:
                            shape.text = title
                            break

            if content:
                from pptx.enum.shapes import PP_PLACEHOLDER
                body_types = {PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT}
                for shape in new_slide.shapes:
                    if shape.has_text_frame and shape.is_placeholder:
                        if shape.placeholder_format.type in body_types:
                            shape.text = content
                            break

            buffer = io.BytesIO()
            try:
                prs.save(buffer)
                buffer.seek(0)
                await overwrite_file_content(context, presentation_id, buffer.read())
            finally:
                buffer.close()

            return ActionResult(
                data={
                    "slide_index": slide_index,
                    "slide_count": len(prs.slides),
                    "result": True
                },
                cost_usd=0.0
            )

        except ImportError:
            return ActionResult(
                data={"result": False, "error": "python-pptx library is required. Install with: pip install python-pptx"},
                cost_usd=0.0
            )
        except Exception as e:
            return ActionResult(
                data={"result": False, "error": str(e)},
                cost_usd=0.0
            )


@microsoft_powerpoint.action("powerpoint_update_slide")
class UpdateSlideAction(ActionHandler):
    """Update the content of an existing slide."""

    async def execute(self, inputs: Dict[str, Any], context: ExecutionContext) -> ActionResult:
        try:
            from pptx import Presentation

            presentation_id = inputs['presentation_id']
            slide_index = inputs['slide_index']
            title = inputs.get('title')
            content = inputs.get('content')
            notes = inputs.get('notes')

            if slide_index < 1:
                return ActionResult(
                    data={"result": False, "error": "Slide index must be 1 or greater"},
                    cost_usd=0.0
                )

            content_bytes = await download_file_content(context, presentation_id)
            prs = Presentation(io.BytesIO(content_bytes))

            if slide_index > len(prs.slides):
                return ActionResult(
                    data={"result": False, "error": f"Slide index {slide_index} is out of range. Presentation has {len(prs.slides)} slides."},
                    cost_usd=0.0
                )

            slide = prs.slides[slide_index - 1]
            updated = False

            if title is not None:
                from pptx.enum.shapes import PP_PLACEHOLDER
                title_types = {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE}
                for shape in slide.shapes:
                    if shape.has_text_frame and shape.is_placeholder:
                        if hasattr(shape, 'placeholder_format') and shape.placeholder_format.type in title_types:
                            shape.text = title
                            updated = True
                            break

            if content is not None:
                from pptx.enum.shapes import PP_PLACEHOLDER
                body_types = {PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT}
                for shape in slide.shapes:
                    if shape.has_text_frame and shape.is_placeholder:
                        if hasattr(shape, 'placeholder_format') and shape.placeholder_format.type in body_types:
                            shape.text = content
                            updated = True
                            break

            if notes is not None:
                try:
                    # Accessing notes_slide may raise if notes don't exist
                    # python-pptx will create the notes slide if it doesn't exist
                    notes_slide = slide.notes_slide
                    notes_slide.notes_text_frame.text = notes
                    updated = True
                except Exception as notes_error:
                    # If notes update fails, continue but report partial success
                    return ActionResult(
                        data={
                            "result": False,
                            "error": f"Failed to update speaker notes: {str(notes_error)}. Title/content may have been updated."
                        },
                        cost_usd=0.0
                    )

            if updated:
                buffer = io.BytesIO()
                try:
                    prs.save(buffer)
                    buffer.seek(0)
                    await overwrite_file_content(context, presentation_id, buffer.read())
                finally:
                    buffer.close()

            return ActionResult(
                data={
                    "updated": updated,
                    "slide_index": slide_index,
                    "result": True
                },
                cost_usd=0.0
            )

        except ImportError:
            return ActionResult(
                data={"result": False, "error": "python-pptx library is required. Install with: pip install python-pptx"},
                cost_usd=0.0
            )
        except Exception as e:
            return ActionResult(
                data={"result": False, "error": str(e)},
                cost_usd=0.0
            )


@microsoft_powerpoint.action("powerpoint_delete_slide")
class DeleteSlideAction(ActionHandler):
    """Delete a slide from the presentation."""

    async def execute(self, inputs: Dict[str, Any], context: ExecutionContext) -> ActionResult:
        try:
            from pptx import Presentation

            presentation_id = inputs['presentation_id']
            slide_index = inputs['slide_index']

            if slide_index < 1:
                return ActionResult(
                    data={"result": False, "error": "Slide index must be 1 or greater"},
                    cost_usd=0.0
                )

            content_bytes = await download_file_content(context, presentation_id)
            prs = Presentation(io.BytesIO(content_bytes))

            if slide_index > len(prs.slides):
                return ActionResult(
                    data={"result": False, "error": f"Slide index {slide_index} is out of range. Presentation has {len(prs.slides)} slides."},
                    cost_usd=0.0
                )

            if len(prs.slides) == 1:
                return ActionResult(
                    data={"result": False, "error": "Cannot delete the last slide. Presentation must have at least one slide."},
                    cost_usd=0.0
                )

            sldIdLst = prs.slides._sldIdLst
            sldId = sldIdLst[slide_index - 1]
            rId = sldId.rId
            prs.part.drop_rel(rId)
            sldIdLst.remove(sldId)

            remaining_count = len(prs.slides)

            buffer = io.BytesIO()
            try:
                prs.save(buffer)
                buffer.seek(0)
                await overwrite_file_content(context, presentation_id, buffer.read())
            finally:
                buffer.close()

            return ActionResult(
                data={
                    "deleted": True,
                    "slide_count": remaining_count,
                    "result": True
                },
                cost_usd=0.0
            )

        except ImportError:
            return ActionResult(
                data={"result": False, "error": "python-pptx library is required. Install with: pip install python-pptx"},
                cost_usd=0.0
            )
        except Exception as e:
            return ActionResult(
                data={"result": False, "error": str(e)},
                cost_usd=0.0
            )


@microsoft_powerpoint.action("powerpoint_export_pdf")
class ExportPdfAction(ActionHandler):
    """Export the presentation to PDF format."""

    async def execute(self, inputs: Dict[str, Any], context: ExecutionContext) -> ActionResult:
        try:
            presentation_id = inputs['presentation_id']
            output_folder = inputs.get('output_folder')
            output_name = inputs.get('output_name')

            file_info = await context.fetch(
                f"{GRAPH_API_BASE_URL}/me/drive/items/{presentation_id}",
                method="GET",
                params={"$select": "name,parentReference"}
            )

            original_name = file_info.get('name', 'presentation.pptx')
            parent_path = file_info.get('parentReference', {}).get('path', '/drive/root:').replace('/drive/root:', '')

            if output_name:
                pdf_name = output_name if output_name.endswith('.pdf') else f"{output_name}.pdf"
            else:
                pdf_name = original_name.replace('.pptx', '.pdf')

            dest_folder = output_folder if output_folder else parent_path

            pdf_content = await context.fetch(
                f"{GRAPH_API_BASE_URL}/me/drive/items/{presentation_id}/content?format=pdf",
                method="GET",
                raw_response=True
            )

            if dest_folder and dest_folder != "/":
                upload_url = f"{GRAPH_API_BASE_URL}/me/drive/root:/{dest_folder.strip('/')}/{pdf_name}:/content"
            else:
                upload_url = f"{GRAPH_API_BASE_URL}/me/drive/root:/{pdf_name}:/content"

            response = await context.fetch(
                upload_url,
                method="PUT",
                headers={"Content-Type": "application/pdf"},
                data=pdf_content
            )

            # Check if upload response already contains download URL
            download_url = response.get('@microsoft.graph.downloadUrl')
            if not download_url:
                # Only fetch again if download URL not in upload response
                download_response = await context.fetch(
                    f"{GRAPH_API_BASE_URL}/me/drive/items/{response.get('id')}",
                    method="GET"
                )
                download_url = download_response.get('@microsoft.graph.downloadUrl')

            return ActionResult(
                data={
                    "pdf_id": response.get('id'),
                    "pdf_name": response.get('name'),
                    "pdf_webUrl": response.get('webUrl'),
                    "pdf_size": response.get('size'),
                    "download_url": download_url,
                    "result": True
                },
                cost_usd=0.0
            )

        except Exception as e:
            return ActionResult(
                data={"result": False, "error": str(e)},
                cost_usd=0.0
            )


@microsoft_powerpoint.action("powerpoint_get_slide_image")
class GetSlideImageAction(ActionHandler):
    """Get a presentation thumbnail image. Note: Microsoft Graph returns a single thumbnail for the entire presentation, not per-slide images."""

    async def execute(self, inputs: Dict[str, Any], context: ExecutionContext) -> ActionResult:
        try:
            presentation_id = inputs['presentation_id']
            slide_index = inputs['slide_index']
            size = inputs.get('size', 'large')
            image_format = inputs.get('format', 'png')

            if slide_index < 1:
                return ActionResult(
                    data={"result": False, "error": "Slide index must be 1 or greater"},
                    cost_usd=0.0
                )

            size_info = THUMBNAIL_SIZES.get(size, THUMBNAIL_SIZES['large'])

            # Fetch thumbnail directly without downloading the entire file
            # Note: Graph API doesn't support per-slide thumbnails, so we just return
            # the presentation thumbnail with a warning
            thumbnails_response = await context.fetch(
                f"{GRAPH_API_BASE_URL}/me/drive/items/{presentation_id}/thumbnails",
                method="GET"
            )

            thumbnail_sets = thumbnails_response.get('value', [])

            if not thumbnail_sets:
                return ActionResult(
                    data={"result": False, "error": "No thumbnail available for this presentation"},
                    cost_usd=0.0
                )

            size_data = thumbnail_sets[0].get(size, thumbnail_sets[0].get('large', {}))

            return ActionResult(
                data={
                    "image_url": size_data.get('url'),
                    "width": size_data.get('width', size_info['width']),
                    "height": size_data.get('height', size_info['height']),
                    "format": image_format,
                    "requested_slide_index": slide_index,
                    "note": "Microsoft Graph provides a single presentation thumbnail, not per-slide images. The slide_index parameter is accepted but the returned image is for the entire presentation.",
                    "result": True
                },
                cost_usd=0.0
            )

        except Exception as e:
            return ActionResult(
                data={"result": False, "error": str(e)},
                cost_usd=0.0
            )
